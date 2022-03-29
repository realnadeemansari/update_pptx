import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum import shapes
from pptx.dml.color import ColorFormat, RGBColor
from pptx.chart.data import ChartData
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def calc_data():    
    df_wpp = pd.read_csv("reports/amit_Weekly_Portfolio_Performance.csv")
    df_hd = pd.read_csv("reports/Amit_Holding_Distribution_by_Indus.csv")
    df_mf = pd.read_excel("reports/Annual Returns.xlsx")
    df_ri = pd.read_csv("reports/Amit_Portfolio_Risk_Ind.csv")
    
    #YTD RETURNS BY PRODUCT		
    algot = (df_wpp["Week End URP"].iloc[-1] - df_wpp["Week Start URP"][0]) / df_wpp["Week Start URP"][0]
    nifty = (df_wpp["NIFTY End"].iloc[-1] - df_wpp["NIFTY Start"][0]) / df_wpp["NIFTY Start"][0]
    df_mf["YTD"] = df_mf["Unnamed: 4"].apply(lambda x: pd.to_numeric(x, errors="coerce"))
    mfe = df_mf["YTD"].dropna().sum() / len(df_mf["YTD"].dropna())
    YTD = (algot, nifty, mfe)
    
    # WEEKLY RETURNS: ALGO T VS NIFTY
    lc_dates = df_wpp["Unnamed: 0"].to_list().copy()
    lc_dates =  [datetime.strptime(date, "%Y-%m-%d") for date in lc_dates]
    df_wpp["portfolio_return"] = df_wpp[["Week Start URP", "Week End URP"]].apply(lambda x: (x[1]/x[0])-1, axis=1)
    df_wpp["nifty_return"] = df_wpp[["NIFTY Start", "NIFTY End"]].apply(lambda x: (x[1]/x[0])-1, axis=1)
    WRAVN = df_wpp[["portfolio_return", "nifty_return"]]
    
    # OVERALL PORTFOLIO VALUE
    portfolio = (algot * 100) + 100
    nifty = (nifty * 100) + 100
    OPV = (portfolio, nifty)
    
    # PORTFOLIO ALLOCATION BY TOP 5 SECTORS
    df_hd = df_hd[df_hd["Unnamed: 0"] == "Average"][["Unnamed: 1", "Values"]]
    df_hd["Unnamed: 1"] = df_hd["Unnamed: 1"].apply(lambda label: label.replace("_", " "))
    df_hd["Values"] = df_hd["Values"].apply(lambda value: int(value) / 100)
    
    # RISK INDICATOR
    avg_hold = df_ri[df_ri["Unnamed: 1"] == "Avg_Hold"]["0"].to_list()[0]
    avg_hold = str(avg_hold)
    start_date = df_ri[df_ri["Unnamed: 1"] == "Start Date"]["0"].to_list()[0]
    week_start = df_ri[df_ri["Unnamed: 1"] == "Week Start"]["0"].to_list()[0]
    week_end = df_ri[df_ri["Unnamed: 1"] == "Week End"]["0"].to_list()[0]
    start_date = datetime.strptime(start_date, "%Y-%m-%d").strftime('%d-%b, %Y')
    week_start = datetime.strptime(week_start, "%Y-%m-%d").strftime('%d-%b, %Y')
    week_end = datetime.strptime(week_end, "%Y-%m-%d").strftime('%d-%b, %Y')
    dates = [start_date, week_start, week_end]
    
    return YTD, WRAVN, OPV, df_hd, dates, lc_dates, avg_hold
    
def generate_pptx():
    i = 0
    j = 0
    try:
        result = calc_data()
        prs = Presentation("pptx/AlgoT_Pitch_17032022.pptx")
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        if j == 16:
                            paragraph.text = "(" + result[4][0] + " – " + result[4][2] + ")"
                            paragraph.runs[0].font.color.rgb = RGBColor(0, 112, 192)
                            paragraph.runs[0].font.size = Pt(8)
                            paragraph.runs[0].font.italic = True
                        elif j == 18:
                            paragraph.text = "(" + result[4][0] + " – " + result[4][2] + ")"
                            paragraph.runs[0].font.color.rgb = RGBColor(0, 112, 192)
                            paragraph.runs[0].font.size = Pt(8)
                            paragraph.runs[0].font.italic = True
                        elif j == 20:
                            paragraph.text = "(" + result[4][1] + " – " + result[4][2] + ")"
                            paragraph.runs[0].font.color.rgb = RGBColor(0, 112, 192)
                            paragraph.runs[0].font.size = Pt(8)
                            paragraph.runs[0].font.italic = True
                        j+=1
                except Exception as e:
                    if i == 0:
                        # UPDATE VERTICAL BAR CHART
                        chart_data = ChartData()
                        chart_data.categories = shape.chart.plots[0].categories
                        chart_data.add_series('New Series 1', result[0], number_format="#.#%")
                        shape.chart.replace_data(chart_data)

                    elif i == 1:
                        # UPDATE LINE CHART
                        chart_data = ChartData()
                        chart_data.categories = result[5]
                        chart_data.add_series('AlgoT', tuple(result[1]["portfolio_return"].to_list()), number_format="#.##%")
                        chart_data.add_series('NIFTY', tuple(result[1]["nifty_return"].to_list()), number_format="#.##%")
                        shape.chart.replace_data(chart_data)

                    elif i == 2:
                        # UPDATE TABLE OVERALL PORTFOLIO VALUE
                        cells = []
                        for table in shape.table.iter_cells():
                            cells.append(table)

                        cells[6].text = "CURRENT VALUE (" + result[4][2] + ")"
                        cells[7].text = "{:.2f}".format(result[2][0])
                        cells[8].text = "{:.2f}".format(result[2][1])
                        cells[6].text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                        cells[6].text_frame.paragraphs[0].runs[0].font.bold = True
                        cells[6].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                        cells[7].text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                        cells[7].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        cells[8].text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                        cells[8].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                    elif i == 6:
                        # UPDATE HORIZONTAL LINE CHART
                        chart_data = ChartData()
                        chart_data.categories = result[3]["Unnamed: 1"].to_list()
                        chart_data.add_series('New Series 1', result[3]["Values"], number_format="##%")
                        shape.chart.replace_data(chart_data)

                    elif i == 7:
                        # UPDATE TABLE AVERAGE HOLDINGS IN DAYS
                        cells = []
                        for table in shape.table.iter_cells():
                            cells.append(table)

                        cells[0].text = 'Average Holding in Days\n(' + result[4][1] + ' to ' + result[4][2] + ')'
                        cells[1].text = result[6]
                        cells[0].text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                        cells[0].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                        cells[0].text_frame.paragraphs[0].runs[0].font.bold = True
                        cells[0].text_frame.paragraphs[1].runs[0].font.size = Pt(6)
                        cells[0].text_frame.paragraphs[1].runs[0].font.bold = False
                        cells[1].text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                        cells[1].text_frame.paragraphs[0].runs[0].font.bold = False
                        cells[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        cells[1].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        cells[0].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        cells[0].text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        
                    i+=1
                    pass
        prs.save("pptx/result_{}.pptx".format(result[4][2]))
        print("Successfully updated pptx")
    except Exception as e:
        print(e)
    
if __name__ == "__main__":
    generate_pptx()